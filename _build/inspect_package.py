from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_docx(path: Path) -> str:
    doc = Document(path)
    lines: list[str] = [f"# {path.name}"]
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(f"P[{paragraph.style.name}]: {text}")
    for index, table in enumerate(doc.tables, start=1):
        lines.append(f"TABLE {index}")
        for row in table.rows:
            lines.append(" | ".join(cell.text.replace("\n", " / ").strip() for cell in row.cells))
    return "\n".join(lines)


def extract_xlsx(path: Path) -> str:
    wb = load_workbook(path, data_only=False, read_only=True)
    lines: list[str] = [f"# {path.name}", "Sheets: " + ", ".join(wb.sheetnames)]
    for ws in wb.worksheets:
        max_row = ws.max_row or 0
        max_column = ws.max_column or 0
        lines.append(f"\n## {ws.title} ({max_row} rows x {max_column} cols)")
        if ws.title == "01_Input":
            for row in ws.iter_rows(min_row=1, max_row=min(max_row, 45), values_only=True):
                values = ["" if value is None else str(value) for value in row[:10]]
                if any(values):
                    lines.append(" | ".join(values))
        else:
            for row in ws.iter_rows(min_row=1, max_row=min(max_row, 20), values_only=True):
                values = ["" if value is None else str(value) for value in row[:8]]
                if any(values):
                    lines.append(" | ".join(values))
    return "\n".join(lines)


def main() -> int:
    package = Path(__file__).resolve().parents[2]
    output_dir = Path(__file__).resolve().parent / "source_extracts"
    output_dir.mkdir(parents=True, exist_ok=True)

    for path in sorted(package.glob("*.docx")):
        (output_dir / f"{path.stem}.txt").write_text(extract_docx(path), encoding="utf-8")
    for name in ("PG_FinancialStory_Student.xlsx", "PG_FinancialStory_Instructor.xlsx"):
        path = package / name
        (output_dir / f"{path.stem}.txt").write_text(extract_xlsx(path), encoding="utf-8")

    course_root = package.parent
    deck_names = [
        "Chapter 1.pptx",
        "Ch2.pptx",
        "Chapter 3.pptx",
        "Chapter 4.pptx",
        "Ch 5.pptx",
        "Ch 6.pptx",
        "Ch 7.pptx",
        "Ch 8.pptx",
    ]
    chapter_lines: list[str] = []
    for name in deck_names:
        deck_path = course_root / name
        prs = Presentation(deck_path)
        chapter_lines.append(f"# {name}")
        for slide_no, slide in enumerate(prs.slides, start=1):
            texts = [shape.text.strip().replace("\n", " / ") for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text.strip()]
            if texts:
                chapter_lines.append(f"{slide_no}: {' || '.join(texts[:3])}")
        chapter_lines.append("")
    (output_dir / "Chapter_Deck_Titles.txt").write_text("\n".join(chapter_lines), encoding="utf-8")

    print(output_dir)
    return 0


if __name__ == "__main__":
    sys.exit(main())
